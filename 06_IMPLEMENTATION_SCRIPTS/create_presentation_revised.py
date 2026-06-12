from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import datetime

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define KMU color scheme (from AI-Powered design image)
COLOR_PRIMARY_NAVY = RGBColor(26, 58, 82)  # Navy Blue #1A3A52
COLOR_SECONDARY_TEAL = RGBColor(46, 139, 158)  # Teal #2E8B9E
COLOR_ACCENT_ORANGE = RGBColor(255, 159, 28)  # Orange #FF9F1C
COLOR_SUCCESS_GREEN = RGBColor(26, 188, 156)  # Teal Green #1ABC9C
COLOR_ALERT_RED = RGBColor(231, 76, 60)  # Crimson #E74C3C
COLOR_TEXT = RGBColor(51, 51, 51)  # Dark Gray
COLOR_WHITE = RGBColor(255, 255, 255)  # White
COLOR_LIGHT_BG = RGBColor(245, 245, 245)  # Light Gray

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide with KMU branding"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY_NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_ACCENT_ORANGE
        p.alignment = PP_ALIGN.CENTER

    # Footer dengan branding KMU
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = "PT. Kaltim Medika Utama — Departemen Pengadaan Umum dan Jasa"
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_SECONDARY_TEAL
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    # Title bar dengan warna KMU
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_PRIMARY_NAVY
    title_shape.line.color.rgb = COLOR_PRIMARY_NAVY

    # Title text
    title_frame = title_shape.text_frame
    title_frame.margin_top = Inches(0.1)
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.LEFT

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)

    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_PRIMARY_NAVY
    title_shape.line.color.rgb = COLOR_PRIMARY_NAVY

    # Title text
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    for i, point in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4), Inches(5.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    for i, point in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    return slide

# ============================================
# PRESENTATION: DIGITALISASI PENGADAAN KMU
# DENGAN NOMENKLATUR RESMI KMU
# ============================================

# Slide 1: Title Slide
add_title_slide(prs, "Digitalisasi Pengadaan KMU",
                "Platform Terintegrasi dengan AI Penjaga Kepatuhan SPO")

# Slide 2: Executive Summary (KMU Terminology)
add_content_slide(prs, "Ringkasan Eksekutif", [
    "Tujuan: Transformasi digital sistem Daan Umum & Daan Jasa PT. KMU",
    "Cakupan: 9 fase siklus pengadaan + dukungan lintas departemen",
    "Inovasi: AI Penjaga Kepatuhan SPO real-time sesuai Pedoman Pengadaan KMU",
    "Stakeholder: Direksi, Dept. Pengadaan, SBU, Vendor, Komite Anggaran, SPI",
    "Target: Efisiensi 30%, Kepatuhan SPO 100%, Transparansi penuh",
    "Status: Siap untuk implementasi IT setelah Direksi approval"
])

# Slide 3: Tantangan Saat Ini
add_content_slide(prs, "Tantangan Saat Ini", [
    "Proses Daan Umum & Daan Jasa tersebar di berbagai sistem manual/semi-digital",
    "Dokumentasi SPO & Pedoman tidak terintegrasi dengan proses operasional",
    "Sulit melacak status PP, SPPH, SJPH, PO, SPK, PKS secara real-time dari Direksi",
    "Risiko ketidakpatuhan SPO tidak terdeteksi sampai audit (terlambat)",
    "Laporan analitik & insights untuk Komite Anggaran memerlukan waktu lama",
    "Vendor Management & Database tidak terintegrasi dalam satu platform"
])

# Slide 4: Solusi Overview
add_content_slide(prs, "Solusi: Platform Terintegrasi", [
    "Satu platform digital untuk seluruh siklus pengadaan 9 fase (sesuai Glosarium KMU)",
    "AI Guardian: Real-time validasi SPO di setiap keputusan (approval matrix KMU)",
    "Command Center: Dashboard real-time untuk monitoring & decision-making Direksi",
    "Vendor Portal: Self-service untuk registrasi, KSO reporting, dan obligations tracking",
    "Integrated Reporting: Analytics & compliance insights otomatis untuk SPI",
    "Audit Trail: Complete traceability untuk semua transaksi (BAPB, Invoice, Payment)"
])

# Slide 5: Arsitektur Sistem (5 Lapisan)
add_content_slide(prs, "Arsitektur Sistem: 5 Lapisan KMU", [
    "Lapisan 1 - Manajemen & Tata Kelola: E-Procurement, Workflow, Compliance Engine",
    "Lapisan 2 - Operasional: Vendor Management, Catalog, Contract Lifecycle (PKS, SPK)",
    "Lapisan 3 - Integrasi: SIMRS, Finance/ERP, Warehouse, Komite Anggaran",
    "Lapisan 4 - Analitik & Transparansi: Dashboard, Spend Analysis, Audit Trail untuk SPI",
    "Lapisan 5 - Infrastruktur & Keamanan: Cloud/Hybrid, SSO, Data Security",
    "Terintegrasi: GLOSARIUM_ISTILAH_KMU, Pedoman Pengadaan, SPO Latest"
])

# Slide 6: Komponen & Modul Utama
add_two_column_slide(prs, "Komponen & Modul Pengadaan KMU",
    [
        "DEPARTEMEN PENGADAAN (Internal):",
        "AI Compliance Guardian (SPO Enforcement)",
        "Command Center Dashboard",
        "Database Vendor Terverifikasi KMU",
        "",
        "SBU & UNIT KERJA (Internal):",
        "Fase 1: Perencanaan Anggaran (RKAP, RAB)",
        "Fase 3: Permintaan Pengadaan (PP, IMT, SPPJ)",
        "Fase 6: Penerimaan & Verifikasi (BAPB)"
    ],
    [
        "VENDOR EKSTERNAL:",
        "Fase 2: Registrasi & Kualifikasi",
        "Fase 5: Portal KSO (Obligations Tracking)",
        "",
        "KEUANGAN & DIREKSI:",
        "Fase 7: Pembayaran & Invoice (3-Way Match)",
        "Fase 8: Evaluasi Vendor (Scoring 7 Dimensi)",
        "Fase 9: Pelaporan & Analitik"
    ]
)

# Slide 7: 9 Fase Siklus Pengadaan KMU
add_content_slide(prs, "9 Fase Siklus Pengadaan Terintegrasi", [
    "Fase 1: Perencanaan & Anggaran - RKAP, RAB, Budget Ceiling dari SBU",
    "Fase 2: Registrasi & Database Vendor - Kualifikasi, Legal Checklist, Master Database",
    "Fase 3: Permintaan Pengadaan - PP (barang umum), IMT (investasi), SPPJ (jasa)",
    "Fase 4: Pelaksanaan Pengadaan - SPPH→SJPH→Tender/Penunjukan Langsung→PO/SPK/PKS",
    "Fase 5: Kontrak & Portal KSO - Vendor self-report KSO, tracking konsumabel (Lab, Farmasi, BMHP)",
    "Fase 6: Penerimaan & Verifikasi - BAPB, QC, 3-Way Match Check, Konfirmasi selesai",
    "Fase 7: Pembayaran & Invoice - DP & Pelunasan, 3-Way Match (PO-BAPB-Invoice), GL Posting",
    "Fase 8: Evaluasi Vendor - Scoring 7 dimensi, Renewal/Blacklist/Improvement decision",
    "Fase 9: Pelaporan & Analitik - Laporan bulanan/triwulan/tahunan, KPI insights"
])

# Slide 8: AI Guardian System (Kepatuhan SPO)
add_content_slide(prs, "AI Penjaga Kepatuhan SPO", [
    "Real-time Validasi: Setiap action validasi terhadap SPO & Pedoman Pengadaan KMU",
    "Document Ingestion: Upload & extraction SPO otomatis dari berbagai format",
    "Procedure Validation Engine: Checking kepatuhan prosedur step-by-step vs SPO",
    "Chatbot Q&A: Guidance system untuk pertanyaan SPO kapanpun (natural language)",
    "Compliance Alerts: Alert instant jika ada penyimpangan prosedur dari SPO",
    "Compliance Metrics: Dashboard menunjukkan skor kepatuhan SPO per departemen",
    "ROI: Estimasi benefit Rp 208 Juta untuk Year 1 (error prevention, audit efficiency)"
])

# Slide 9: Command Center Dashboard
add_content_slide(prs, "Command Center Dashboard (untuk Direksi)", [
    "Single Screen untuk Direksi: Monitor SEMUA aspek pengadaan real-time",
    "9 Section Dashboard: KPI, Tender Status, Alerts, Payments, Vendor Performance",
    "SLA Countdown Timers: Tracking deadline untuk setiap fase pengadaan (2h, 14h, 7h, 30h)",
    "Alert Center: Escalation otomatis untuk issues & risks (vendor delays, approval bottleneck)",
    "Budget Tracking: Real-time budget consumption vs RKAP & forecasting",
    "Vendor Performance: Rating vendor per 7 dimensi (delivery, quality, harga, etc)",
    "Live Audit Trail: Setiap event tercatat untuk SPI & external audit (BPK)"
])

# Slide 10: Fitur Vendor Management
add_content_slide(prs, "Vendor Management & Portal KSO", [
    "Self-Service Registration: Vendor dapat register & upload dokumen langsung ke sistem",
    "Automated Legal Checklist: Validasi dokumen legal otomatis vs Pedoman KMU",
    "Master Database: Vendor database terpusat dengan status (Aktif, Blacklist, etc)",
    "KSO Reporting: Vendor self-report konsumabel, maintenance, obligations per contract",
    "Performance Tracking: 7 dimensi scoring untuk evaluasi vendor (Delivery, Quality, Price, Compliance, Financial, Response, Partnership)",
    "Automated Notifications: Status updates & requirements push ke vendor via sistem",
    "Integration dengan Procurement: Data vendor siap di procurement process (Fase 3-8)"
])

# Slide 11: Integrasi Multi-Stakeholder
add_two_column_slide(prs, "Integrasi Multi-Stakeholder KMU",
    [
        "DEPARTEMEN PENGADAAN:",
        "Manage vendor database",
        "Process tender & kontrak (PO, SPK, PKS)",
        "Monitor compliance SPO real-time",
        "Generate reports untuk Direksi",
        "",
        "SBU & UNIT KERJA:",
        "Submit kebutuhan & budget (RKAP)",
        "Submit permintaan pengadaan (PP, IMT, SPPJ)",
        "Terima barang/jasa",
        "Verify penerimaan (BAPB sign-off)"
    ],
    [
        "VENDOR EKSTERNAL:",
        "Register & upload dokumen",
        "Self-report KSO obligations",
        "Submit invoices (Invoice)",
        "Track payments & delivery status",
        "",
        "KEUANGAN & DIREKSI:",
        "Authorize budget & payment (Komite Anggaran)",
        "Review spending patterns & ROI",
        "Monitor cash flow & working capital",
        "Strategic decision making (Direksi)"
    ]
)

# Slide 12: Teknologi & Technical Stack
add_content_slide(prs, "Teknologi & Technical Stack", [
    "Database: PostgreSQL - relational database untuk data integrity & audit trail",
    "Backend: Python - API services, AI processing, business logic untuk validasi",
    "Frontend: React - responsive UI untuk desktop & mobile (Web + Mobile App)",
    "Chatbot Widget: React + Gemini Chat API untuk guidance interface",
    "AI/ML: Google Gemini - procedure validation & compliance checking engine",
    "Security: JWT authentication, role-based access control (sesuai otorisasi KMU), encryption",
    "Real-time: WebSocket untuk live updates & notifications ke dashboard"
])

# Slide 13: Manfaat & ROI
add_content_slide(prs, "Manfaat & ROI (Year 1 & Beyond)", [
    "Efisiensi: Proses pengadaan 30% lebih cepat dengan workflow otomatis (5 hari saved per tender)",
    "Kepatuhan: 100% compliance terhadap SPO dengan AI guardian real-time",
    "Cost Savings: Negosiasi lebih baik, vendor performance tracking, reduce waste (Rp 4.2B/year)",
    "Transparansi: Visibility penuh untuk Direksi, audit trail lengkap untuk SPI & BPK",
    "Quality: Standardized process, quality metrics tracking, vendor scoring sistematis",
    "Intelligence: Data-driven insights untuk decision making strategis (spend analysis, trends)",
    "Satisfaction: Faster service untuk SBU, transparency untuk vendor, compliance untuk audit"
])

# Slide 14: Implementation Timeline
add_content_slide(prs, "Implementation Timeline (12 Minggu)", [
    "Week 1-2: Setup & Validasi (Rp 150M) - POC Gemini API, infrastructure",
    "Week 3-6: Core AI Development (Rp 600M) - Document processor, validator, knowledge base",
    "Week 7-8: Chatbot & Integration (Rp 150M) - Guidance system, widget development",
    "Week 9-10: Platform Integration (Rp 100M) - Connect ke procurement platform, E2E testing",
    "Week 11: Testing & Optimization (Rp 50M) - Load testing, security audit, UAT",
    "Week 12: Training & Launch (Rp 50M) - User training, soft launch, production rollout",
    "Total Investment: Rp 1.04 Billion (one-time) + Rp 600M/year (maintenance)"
])

# Slide 15: Success Metrics & KPI
add_two_column_slide(prs, "Success Metrics & KPI Targets",
    [
        "OPERATIONAL:",
        "Proses time: Reduced 30%",
        "Cycle time: Reduced 40%",
        "Error rate: Reduced 95%",
        "Manual work: Reduced 75%",
        "",
        "COMPLIANCE:",
        "SPO adherence: >= 99%",
        "Audit findings: Reduced 90%"
    ],
    [
        "FINANCIAL:",
        "Cost saving: Rp 500M+/year",
        "Payment on-time: >= 98%",
        "Working capital: Improved 20%",
        "",
        "SATISFACTION:",
        "SBU satisfaction: >= 4.5/5",
        "Vendor satisfaction: >= 4/5",
        "SPI audit efficiency: +60%"
    ]
)

# Slide 16: Risk Management
add_content_slide(prs, "Risk Management & Mitigation", [
    "Risk 1: User Adoption - Mitigation: Comprehensive training & change management (4h per user)",
    "Risk 2: Data Migration - Mitigation: Careful planning & validation dengan QA",
    "Risk 3: System Performance - Mitigation: Load testing & optimization (<500ms SLA)",
    "Risk 4: SPO Quality - Mitigation: 3-layer review dari Dept. Pengadaan sebelum digitize",
    "Risk 5: Integration Issues - Mitigation: Pre-testing dengan existing systems (SIMRS, Finance)",
    "Risk 6: Vendor Readiness - Mitigation: Gradual onboarding & portal support untuk vendor"
])

# Slide 17: Next Steps
add_content_slide(prs, "Langkah Selanjutnya", [
    "Approval dari Direksi untuk proceed dengan development (formal SKD)",
    "Setup IT infrastructure & environment (GCP, databases, APIs)",
    "Finalize technical specifications & database schema dengan Dept. Pengadaan",
    "Identify & train implementation team (13 FTE untuk 12 minggu)",
    "Prepare user documentation & training materials (sesuai GLOSARIUM)",
    "Schedule kickoff meeting dengan seluruh stakeholder (Dept. Pengadaan, SBU, Keuangan, SPI)",
    "Begin Phase 1: System Setup & Database (setelah approval)"
])

# Slide 18: Closing
add_title_slide(prs, "Terima Kasih",
                "Mari Transformasikan Pengadaan KMU dengan AI & Kepatuhan SPO")

# Save presentation dengan naming baru
output_path = r"C:\Users\Petugas\Claude\Projects\Pengadaan\Presentasi_Digitalisasi_Pengadaan_KMU_REVISED.pptx"
prs.save(output_path)
print("[OK] Revised Presentation created successfully!")
print(f"[PATH] Saved to: {output_path}")
print(f"[SLIDES] Total slides: {len(prs.slides)}")
print(f"[COLORS] Using KMU brand colors: Navy, Teal, Orange")
print(f"[TERMINOLOGY] All content uses GLOSARIUM_ISTILAH_KMU (Daan Umum, Daan Jasa, SPPH, SJPH, PO, SPK, PKS, BAPB, etc)")
