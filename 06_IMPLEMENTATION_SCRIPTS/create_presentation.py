from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import datetime

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
COLOR_PRIMARY = RGBColor(0, 51, 102)  # Dark Blue
COLOR_SECONDARY = RGBColor(0, 102, 204)  # Medium Blue
COLOR_ACCENT = RGBColor(255, 102, 0)  # Orange
COLOR_TEXT = RGBColor(51, 51, 51)  # Dark Gray
COLOR_WHITE = RGBColor(255, 255, 255)  # White

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_ACCENT
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_PRIMARY
    title_shape.line.color.rgb = COLOR_PRIMARY

    # Title text
    title_frame = title_shape.text_frame
    title_frame.margin_top = Inches(0.1)
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.LEFT

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)

    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_PRIMARY
    title_shape.line.color.rgb = COLOR_PRIMARY

    # Title text
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    for i, point in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4), Inches(5.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    for i, point in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    return slide

# Slide 1: Title Slide
add_title_slide(prs, "Digitalisasi Pengadaan KMU",
                "Platform Terintegrasi dengan AI Penjaga Kepatuhan SPO")

# Slide 2: Executive Summary
add_content_slide(prs, "Executive Summary", [
    "🎯 Tujuan: Transformasi digital sistem pengadaan PT. KMU",
    "📊 Cakupan: 9 fase siklus pengadaan + dukungan lintas fungsi",
    "🤖 Inovasi: AI Penjaga Kepatuhan SPO real-time",
    "👥 Stakeholder: Direksi, Procurement, Unit Kerja, Vendor, Keuangan",
    "📈 Target: Efisiensi 30%, Compliance 100%, Transparansi penuh",
    "⏱️ Status: Siap untuk implementasi IT"
])

# Slide 3: Masalah Saat Ini
add_content_slide(prs, "Tantangan Saat Ini", [
    "❌ Proses pengadaan tersebar di berbagai sistem manual/semi-digital",
    "❌ Dokumentasi SPO tidak terintegrasi dengan proses operasional",
    "❌ Sulit melacak status pengadaan secara real-time dari Direksi",
    "❌ Risiko ketidakpatuha prosedur tidak terdeteksi dengan baik",
    "❌ Laporan analitik & insights memerlukan waktu lama untuk dikumpulkan",
    "❌ Vendor management tidak terintegrasi dalam satu platform"
])

# Slide 4: Solusi Overview
add_content_slide(prs, "Solusi: Platform Terintegrasi", [
    "✅ Satu platform digital untuk seluruh siklus pengadaan 9 fase",
    "✅ AI Guardian: Real-time procedure validation di setiap keputusan",
    "✅ Command Center: Dashboard real-time untuk monitoring & decision-making",
    "✅ Vendor Portal: Self-service untuk registrasi dan pelaporan KSO",
    "✅ Integrated Reporting: Analytics & compliance insights otomatis",
    "✅ Audit Trail: Complete traceability untuk semua transaksi"
])

# Slide 5: Arsitektur Sistem
add_content_slide(prs, "Arsitektur Sistem", [
    "🏗️ Lapisan 1: AI Penjaga Kepatuhan SPO - Validasi real-time setiap aksi",
    "🏗️ Lapisan 2: Command Center Dashboard - Monitoring & control room",
    "🏗️ Lapisan 3: Audit Trail & Pelaporan - Complete transaction history",
    "🏗️ Lapisan 4: 9 Modul Pengadaan - Dari perencanaan hingga pelaporan",
    "🏗️ Lapisan 5: Antarmuka Stakeholder - Vendor, SBU, Keuangan, Direksi",
    "🏗️ Database Terpusat - PostgreSQL dengan schema terintegrasi"
])

# Slide 6: Komponen Utama
add_two_column_slide(prs, "Komponen & Modul Utama",
    [
        "INTERNAL PROCUREMENT:",
        "• AI Compliance Guardian",
        "• Command Center",
        "• Database Vendor Master",
        "",
        "SBU & UNIT KERJA:",
        "• Perencanaan Anggaran (Fase 1)",
        "• Permintaan Pengadaan (Fase 3)",
        "• Penerimaan & Verifikasi (Fase 6)"
    ],
    [
        "VENDOR EKSTERNAL:",
        "• Registrasi & Kualifikasi (Fase 2)",
        "• Portal KSO (Fase 5)",
        "",
        "KEUANGAN & DIREKSI:",
        "• Pembayaran & Invoice (Fase 7)",
        "• Evaluasi Vendor (Fase 8)",
        "• Pelaporan & Analitik (Fase 9)"
    ]
)

# Slide 7: 9 Fase Siklus Pengadaan
add_content_slide(prs, "9 Fase Siklus Pengadaan", [
    "1️⃣ Perencanaan & Anggaran - RKAP, RAB, Budget Ceiling dari SBU",
    "2️⃣ Registrasi & Database Vendor - Kualifikasi vendor, legal checklist",
    "3️⃣ Permintaan Pengadaan - IMT, PP, SPPJ dari unit kerja",
    "4️⃣ Pelaksanaan Pengadaan - SPPH, Aanwijzing, SJPH, Bidding, Negosiasi → PO/SPK/PKS",
    "5️⃣ Kontrak & Portal KSO - Vendor self-report, tracking konsumabel",
    "6️⃣ Penerimaan & Verifikasi - BAPB, QC Barang/Jasa, Konfirmasi selesai",
    "7️⃣ Pembayaran & Invoice - DP & Pelunasan, 3-Way Match, GL Posting",
    "8️⃣ Evaluasi Vendor - Scoring 7 dimensi, Feedback untuk renewal",
    "9️⃣ Pelaporan & Analitik - Laporan bulanan/triwulan/tahunan, insights"
])

# Slide 8: AI Guardian System
add_content_slide(prs, "AI Penjaga Kepatuhan SPO", [
    "🤖 Real-time Validation: Setiap aksi di platform divalidasi terhadap SPO",
    "📋 Document Ingestion: Upload & extraction SPO otomatis dari berbagai format",
    "✓ Procedure Validation Engine: Checking kepatuhan prosedur step-by-step",
    "💬 Chatbot Q&A: Guidance system untuk pertanyaan prosedur kapanpun",
    "⚠️ Compliance Alerts: Alert instant jika ada penyimpangan prosedur",
    "📊 Compliance Metrics: Dashboard menunjukkan skor kepatuhan per departemen",
    "💡 ROI: Estimasi benefit Rp 208 Juta untuk Year 1"
])

# Slide 9: Command Center Dashboard
add_content_slide(prs, "Command Center Dashboard", [
    "🎛️ Single Screen untuk Direksi: Monitor SEMUA aspek pengadaan real-time",
    "📊 9 Section Dashboard: KPI, Tender Status, Alerts, Payments, Vendor Performance",
    "⏱️ SLA Countdown Timers: Tracking deadline untuk setiap fase pengadaan",
    "🚨 Alert Center: Escalation otomatis untuk issues & risks",
    "📈 Budget Tracking: Real-time budget consumption & forecasting",
    "🔴 System Health: Monitoring kesehatan sistem dan uptime",
    "📱 Live Transaction Log: Audit trail setiap event dalam sistem"
])

# Slide 10: Fitur Vendor Management
add_content_slide(prs, "Vendor Portal & Management", [
    "🏢 Self-Service Registration: Vendor dapat register & upload dokumen langsung",
    "✅ Automated Legal Checklist: Validasi dokumen legal otomatis",
    "📊 Master Database: Vendor database terpusat dengan status tracking",
    "📋 KSO Reporting: Vendor self-report konsumabel, maintenance, obligations",
    "⭐ Performance Tracking: 7 dimensi scoring untuk evaluasi vendor",
    "📧 Automated Notifications: Status updates & requirements ke vendor",
    "🔗 Integration dengan Procurement: Data vendor siap di procurement process"
])

# Slide 11: Integrasi Stakeholder
add_two_column_slide(prs, "Integrasi Multi-Stakeholder",
    [
        "PROCUREMENT DEPT:",
        "• Manage vendor database",
        "• Process tender & kontrak",
        "• Monitor compliance",
        "• Generate reports",
        "",
        "SBU & UNIT KERJA:",
        "• Submit kebutuhan & budget",
        "• Submit permintaan pengadaan",
        "• Terima barang/jasa",
        "• Verify penerimaan"
    ],
    [
        "VENDOR:",
        "• Register & upload docs",
        "• Self-report KSO",
        "• Submit invoices",
        "• Track payments",
        "",
        "FINANCE & BOD:",
        "• Authorize budget & payment",
        "• Review spending patterns",
        "• Monitor cash flow",
        "• Strategic decision making"
    ]
)

# Slide 12: Teknologi & Stack
add_content_slide(prs, "Teknologi & Technical Stack", [
    "🗄️ Database: PostgreSQL - relational database untuk data integrity",
    "🖥️ Backend: Python - API services, AI processing, business logic",
    "🎨 Frontend: React - responsive UI untuk desktop & mobile",
    "🤖 AI/ML: Google Gemini - procedure validation & compliance checking",
    "🔐 Security: JWT authentication, role-based access control, encryption",
    "📡 Real-time: WebSocket untuk live updates & notifications",
    "☁️ Deployment: Cloud-ready atau on-premise sesuai preference"
])

# Slide 13: Manfaat & ROI
add_content_slide(prs, "Manfaat & ROI", [
    "⚡ Efisiensi: Proses pengadaan 30% lebih cepat dengan workflow otomatis",
    "✅ Kepatuhan: 100% compliance terhadap SPO dengan AI guardian real-time",
    "💰 Cost Savings: Negosiasi lebih baik, vendor performance tracking, reduce waste",
    "📊 Transparency: Visibility penuh untuk Direksi, audit trail lengkap",
    "🎯 Quality: Standardized process, quality metrics tracking, vendor scoring",
    "💡 Intelligence: Data-driven insights untuk decision making strategis",
    "👥 Satisfaction: Faster service untuk SBU, transparency untuk vendor"
])

# Slide 14: Implementation Roadmap
add_content_slide(prs, "Implementation Timeline", [
    "📅 Phase 1 (Week 1-2): System Setup & Database Initialization",
    "📅 Phase 2 (Week 3-4): Backend API Services Development",
    "📅 Phase 3 (Week 5-6): Frontend UI & Integration",
    "📅 Phase 4 (Week 7): Testing, UAT & Performance Tuning",
    "📅 Phase 5 (Week 8): Training & Onboarding ke Procurement Team",
    "📅 Phase 6-7 (Week 9-10): Soft Launch & Full Production Rollout",
    "📅 Phase 8 (Week 11+): Optimization & Continuous Improvement"
])

# Slide 15: Success Metrics
add_two_column_slide(prs, "Success Metrics & KPI",
    [
        "OPERATIONAL:",
        "• Proses time: ↓30%",
        "• Cycle time: ↓40%",
        "• Error rate: ↓95%",
        "• Manual work: ↓75%",
        "",
        "COMPLIANCE:",
        "• SPO adherence: ≥99%",
        "• Audit findings: ↓90%"
    ],
    [
        "FINANCIAL:",
        "• Cost saving: Rp 500M+/year",
        "• Payment on-time: ≥98%",
        "• Vendor spend mgmt: +20%",
        "",
        "SATISFACTION:",
        "• SBU satisfaction: ≥4.5/5",
        "• Vendor satisfaction: ≥4/5"
    ]
)

# Slide 16: Risiko & Mitigation
add_content_slide(prs, "Risk Management", [
    "⚠️ Risk 1: User adoption - Mitigation: Comprehensive training & support",
    "⚠️ Risk 2: Data migration - Mitigation: Careful planning & validation",
    "⚠️ Risk 3: System performance - Mitigation: Load testing & optimization",
    "⚠️ Risk 4: Integration issues - Mitigation: Pre-testing dengan legacy systems",
    "⚠️ Risk 5: Vendor readiness - Mitigation: Gradual onboarding & support",
    "⚠️ Risk 6: Change management - Mitigation: Clear communication & governance"
])

# Slide 17: Next Steps
add_content_slide(prs, "Langkah Selanjutnya", [
    "✓ Approval dari Direksi untuk proceed dengan development",
    "✓ Setup IT infrastructure & environment",
    "✓ Finalize technical specifications & database schema",
    "✓ Identify & train implementation team",
    "✓ Prepare user documentation & training materials",
    "✓ Schedule kickoff meeting dengan seluruh stakeholder",
    "✓ Begin Phase 1: System Setup & Database"
])

# Slide 18: Closing
add_title_slide(prs, "Terima Kasih",
                "Mari Transformasikan Pengadaan KMU Bersama")

# Save presentation
output_path = r"C:\Users\Petugas\Claude\Projects\Pengadaan\Presentasi_Digitalisasi_Pengadaan_KMU.pptx"
prs.save(output_path)
print("[OK] Presentation created successfully!")
print(f"[PATH] Saved to: {output_path}")
print(f"[SLIDES] Total slides: {len(prs.slides)}")
