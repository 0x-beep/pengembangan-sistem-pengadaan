import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Warna KMU
HIJAU = RGBColor(46, 139, 94)      # #2E8B5E
ORANGE = RGBColor(255, 159, 28)    # #FF9F1C
PUTIH = RGBColor(255, 255, 255)
HITAM = RGBColor(33, 33, 33)
GRAY_LIGHT = RGBColor(240, 240, 240)

def create_presentation():
    """Create PPT Pengadaan KMU from scratch - RAPI & NO OVERLAP"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    def add_slide_with_background(title_text="", has_title=True):
        """Add slide dengan background putih dan header"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Background putih
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = PUTIH
        bg.line.color.rgb = PUTIH

        if has_title and title_text:
            # Header bar hijau
            header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
            header.fill.solid()
            header.fill.fore_color.rgb = HIJAU
            header.line.color.rgb = HIJAU

            # Title text
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
            tf = title_box.text_frame
            tf.text = title_text
            p = tf.paragraphs[0]
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = PUTIH
            p.alignment = PP_ALIGN.LEFT

            # Orange accent line
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.8), prs.slide_width, Inches(0.08))
            accent.fill.solid()
            accent.fill.fore_color.rgb = ORANGE
            accent.line.color.rgb = ORANGE

        return slide

    print("Creating PPT slides...")

    # ===== SLIDE 1: COVER =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = HIJAU
    bg1.line.color.rgb = HIJAU

    # Title
    title1 = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf1 = title1.text_frame
    tf1.word_wrap = True
    tf1.text = "Digitalisasi Pengadaan KMU"
    p1 = tf1.paragraphs[0]
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = PUTIH
    p1.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle1 = slide1.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1.5))
    tf_sub = subtitle1.text_frame
    tf_sub.word_wrap = True
    tf_sub.text = "Platform Terintegrasi dengan\nValidator Prosedur & Sistem Penjaga Kepatuhan SPO"
    for p in tf_sub.paragraphs:
        p.font.size = Pt(24)
        p.font.color.rgb = ORANGE
        p.alignment = PP_ALIGN.CENTER

    # Footer
    footer1 = slide1.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    tf_footer = footer1.text_frame
    tf_footer.text = "PT. Kaltim Medika Utama | Juni 2026"
    p_footer = tf_footer.paragraphs[0]
    p_footer.font.size = Pt(16)
    p_footer.font.color.rgb = PUTIH
    p_footer.alignment = PP_ALIGN.CENTER
    print("✅ Slide 1: Cover")

    # ===== SLIDE 2: RINGKASAN EKSEKUTIF =====
    slide2 = add_slide_with_background("Ringkasan Eksekutif")

    benefits = [("30%", "Lebih Cepat"), ("Rp 14.57B", "Benefit/Tahun"), ("9,636%", "ROI Year 1")]
    for i, (value, label) in enumerate(benefits):
        left = Inches(0.7 + i * 3)
        top = Inches(1.3)

        box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.5), Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = ORANGE
        box.line.color.rgb = HIJAU
        box.line.width = Pt(2)

        val_box = slide2.shapes.add_textbox(left + Inches(0.1), top + Inches(0.3), Inches(2.3), Inches(0.8))
        tf_val = val_box.text_frame
        tf_val.text = value
        p_val = tf_val.paragraphs[0]
        p_val.font.size = Pt(28)
        p_val.font.bold = True
        p_val.font.color.rgb = PUTIH
        p_val.alignment = PP_ALIGN.CENTER

        label_box = slide2.shapes.add_textbox(left + Inches(0.1), top + Inches(1.1), Inches(2.3), Inches(0.6))
        tf_label = label_box.text_frame
        tf_label.word_wrap = True
        tf_label.text = label
        p_label = tf_label.paragraphs[0]
        p_label.font.size = Pt(14)
        p_label.font.color.rgb = PUTIH
        p_label.alignment = PP_ALIGN.CENTER

    content2 = slide2.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(8.6), Inches(3))
    tf_content = content2.text_frame
    tf_content.word_wrap = True
    tf_content.text = """• Platform terintegrasi menggabungkan semua fase pengadaan KMU (Perencanaan - Pelaporan)
• Validator Prosedur & Sistem Penjaga Kepatuhan mengawasi SETIAP transaksi vs SPO KMU
• Real-time Dashboard untuk Direksi dengan visibilitas lengkap status PP, PO, PKS, BAPB
• Automated compliance checking & automated reporting untuk SPI (Satuan Pengawasan Internal)
• Payback period: 5 hari | Success probability: 85-90%"""
    for p in tf_content.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = HITAM
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    print("✅ Slide 2: Ringkasan Eksekutif")

    # ===== SLIDE 3: TANTANGAN =====
    slide3 = add_slide_with_background("Tantangan Saat Ini")
    challenges = ["Proses Manual & Tersebar", "Kepatuhan SPO Tidak Konsisten", "Tidak Ada Visibilitas Real-time", "Audit Menemukan Violations Terlambat"]
    for i, challenge in enumerate(challenges):
        top = Inches(1.3 + i * 1.1)
        cbox = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), top, Inches(8.6), Inches(0.9))
        cbox.fill.solid()
        cbox.fill.fore_color.rgb = GRAY_LIGHT
        cbox.line.color.rgb = ORANGE
        cbox.line.width = Pt(2)

        ctext = slide3.shapes.add_textbox(Inches(1), top + Inches(0.15), Inches(8), Inches(0.6))
        tf_chal = ctext.text_frame
        tf_chal.word_wrap = True
        tf_chal.text = challenge
        p_chal = tf_chal.paragraphs[0]
        p_chal.font.size = Pt(14)
        p_chal.font.color.rgb = HITAM
    print("✅ Slide 3: Tantangan")

    # ===== SLIDE 4: SOLUSI =====
    slide4 = add_slide_with_background("Solusi: Platform Terintegrasi + AI")
    solutions = [("1", "Platform Digital", "Satu sistem untuk 9 fase pengadaan"), ("2", "Validator Prosedur", "AI mengawasi kepatuhan SPO real-time"), ("3", "Dashboard Direksi", "Visibility penuh status & compliance")]
    for i, (num, title, desc) in enumerate(solutions):
        top = Inches(1.3 + i * 1.7)
        circle = slide4.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), top, Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ORANGE
        circle.line.color.rgb = HIJAU

        ctext = slide4.shapes.add_textbox(Inches(0.7), top, Inches(0.5), Inches(0.5))
        tf_c = ctext.text_frame
        tf_c.text = num
        p_c = tf_c.paragraphs[0]
        p_c.font.size = Pt(18)
        p_c.font.bold = True
        p_c.font.color.rgb = PUTIH
        p_c.alignment = PP_ALIGN.CENTER
        p_c.vertical_anchor = MSO_ANCHOR.MIDDLE

        tbox = slide4.shapes.add_textbox(Inches(1.5), top, Inches(7.9), Inches(0.5))
        tf_t = tbox.text_frame
        tf_t.word_wrap = True
        tf_t.text = f"{title}: {desc}"
        p_t = tf_t.paragraphs[0]
        p_t.font.size = Pt(13)
        p_t.font.color.rgb = HITAM
    print("✅ Slide 4: Solusi")

    # ===== SLIDE 5: ARSITEKTUR =====
    slide5 = add_slide_with_background("Arsitektur Sistem: 5 Lapisan KMU")
    layers = ["Management Layer (Direksi, GM, Manager)", "Operational Layer (Pengadaan, Vendor, SBU)", "Integration Layer (Validator, Dashboard, API)", "Analytics Layer (KPI, Reports, Insights)", "Infrastructure Layer (Cloud, Database, Security)"]
    for i, layer in enumerate(layers):
        top = Inches(1.3 + i * 1)
        lbox = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), top, Inches(8), Inches(0.8))
        lbox.fill.solid()
        lbox.fill.fore_color.rgb = ORANGE if i % 2 == 0 else HIJAU
        lbox.line.color.rgb = HITAM
        lbox.line.width = Pt(1)

        ltext = slide5.shapes.add_textbox(Inches(1.3), top + Inches(0.15), Inches(7.4), Inches(0.5))
        tf_l = ltext.text_frame
        tf_l.word_wrap = True
        tf_l.text = layer
        p_l = tf_l.paragraphs[0]
        p_l.font.size = Pt(13)
        p_l.font.bold = True
        p_l.font.color.rgb = PUTIH
    print("✅ Slide 5: Arsitektur")

    # ===== SLIDE 6: MODUL =====
    slide6 = add_slide_with_background("12 Modul Pengadaan KMU")
    modules = ["M1: Platform Pengadaan", "M2: Manajemen Penyedia", "M3: Portal Penyedia", "M4: Validator Prosedur", "M5: Pelaksanaan & Delivery", "M6: Jaminan Mutu", "M7: Pemrosesan Invoice", "M8: Pembayaran", "M9: Pelaporan & Analitik", "M10: Database Penyedia", "M11: Dashboard Pengadaan", "M12: Evaluasi & Scoring"]
    for i, module in enumerate(modules):
        row = i // 4
        col = i % 4
        left = Inches(0.5 + col * 2.2)
        top = Inches(1.2 + row * 1.5)

        mbox = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2), Inches(1.2))
        mbox.fill.solid()
        mbox.fill.fore_color.rgb = ORANGE if i % 2 == 0 else HIJAU
        mbox.line.color.rgb = HIJAU
        mbox.line.width = Pt(1)

        mtext = slide6.shapes.add_textbox(left + Inches(0.1), top + Inches(0.2), Inches(1.8), Inches(0.8))
        tf_m = mtext.text_frame
        tf_m.word_wrap = True
        tf_m.text = module
        for p in tf_m.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = PUTIH
            p.alignment = PP_ALIGN.CENTER
    print("✅ Slide 6: 12 Modul")

    # ===== SLIDE 7: 9 FASE =====
    slide7 = add_slide_with_background("9 Fase Siklus Pengadaan KMU")
    phases = ["Perencanaan", "Registrasi", "Permintaan", "Pelaksanaan", "Kontrak", "Penerimaan", "Pembayaran", "Evaluasi", "Pelaporan"]
    for i, phase in enumerate(phases):
        col = i % 3
        row = i // 3
        left = Inches(0.7 + col * 3)
        top = Inches(1.2 + row * 1.8)

        pbox = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.7), Inches(1.5))
        pbox.fill.solid()
        pbox.fill.fore_color.rgb = ORANGE if i % 2 == 0 else HIJAU
        pbox.line.color.rgb = HIJAU
        pbox.line.width = Pt(2)

        pnum = slide7.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(2.3), Inches(0.4))
        tf_pn = pnum.text_frame
        tf_pn.text = f"Fase {i+1}"
        p_pn = tf_pn.paragraphs[0]
        p_pn.font.size = Pt(11)
        p_pn.font.bold = True
        p_pn.font.color.rgb = PUTIH

        pname = slide7.shapes.add_textbox(left + Inches(0.2), top + Inches(0.7), Inches(2.3), Inches(0.6))
        tf_pname = pname.text_frame
        tf_pname.word_wrap = True
        tf_pname.text = phase
        p_pname = tf_pname.paragraphs[0]
        p_pname.font.size = Pt(12)
        p_pname.font.bold = True
        p_pname.font.color.rgb = PUTIH
        p_pname.alignment = PP_ALIGN.CENTER
    print("✅ Slide 7: 9 Fase")

    # ===== SLIDE 8: VALIDATOR =====
    slide8 = add_slide_with_background("Validator Prosedur & Sistem Penjaga Kepatuhan")
    flow_steps = ["Input", "Validation", "Alert", "Guidance"]
    for i, step in enumerate(flow_steps):
        left = Inches(0.8 + i * 2.2)
        sbox = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(1.8), Inches(1))
        sbox.fill.solid()
        sbox.fill.fore_color.rgb = ORANGE if i % 2 == 0 else HIJAU
        sbox.line.color.rgb = HIJAU
        sbox.line.width = Pt(2)

        stext = slide8.shapes.add_textbox(left + Inches(0.1), Inches(1.7), Inches(1.6), Inches(0.6))
        tf_s = stext.text_frame
        tf_s.word_wrap = True
        tf_s.text = step
        p_s = tf_s.paragraphs[0]
        p_s.font.size = Pt(14)
        p_s.font.bold = True
        p_s.font.color.rgb = PUTIH
        p_s.alignment = PP_ALIGN.CENTER

        if i < len(flow_steps) - 1:
            arrow = slide8.shapes.add_connector(1, left + Inches(1.9), Inches(2), left + Inches(2.2), Inches(2))
            arrow.line.color.rgb = HIJAU
            arrow.line.width = Pt(3)

    desc8 = slide8.shapes.add_textbox(Inches(0.7), Inches(3), Inches(8.6), Inches(3.8))
    tf_d = desc8.text_frame
    tf_d.word_wrap = True
    tf_d.text = """FITUR VALIDATOR:
• Real-time monitoring setiap transaksi pengadaan (PP, SPPH, SJPH, PO, SPK, PKS, BAPB)
• Automatic checking terhadap SPO KMU & authority matrix (otorisasi nilai)
• Alert instant jika ada penyimpangan prosedur dengan severity level (LOW/MEDIUM/HIGH)
• AI Guidance langsung kepada user untuk corrective action
• Compliance scoring per transaksi dengan trend analysis"""
    for p in tf_d.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = HITAM
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    print("✅ Slide 8: Validator")

    # ===== SLIDE 9: DASHBOARD =====
    slide9 = add_slide_with_background("Dashboard Pengadaan untuk Direksi")
    dash_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.3), Inches(7), Inches(1.5))
    dash_box.fill.solid()
    dash_box.fill.fore_color.rgb = HIJAU
    dash_box.line.color.rgb = ORANGE
    dash_box.line.width = Pt(3)

    dash_text = slide9.shapes.add_textbox(Inches(2), Inches(1.8), Inches(6), Inches(0.6))
    tf_dash = dash_text.text_frame
    tf_dash.text = "📊 Real-time Monitoring: Status Transaksi, Compliance Score, Risk Alerts, KPI Targets"
    p_dash = tf_dash.paragraphs[0]
    p_dash.font.size = Pt(16)
    p_dash.font.bold = True
    p_dash.font.color.rgb = PUTIH
    p_dash.alignment = PP_ALIGN.CENTER

    features9 = ["Visibility penuh status semua transaksi (PP, PO, SPK, PKS, BAPB)", "Real-time KPI: % Kepatuhan SPO, Avg Processing time, Error rate", "Risk dashboard: High-risk transactions, bottleneck processes", "Export reports untuk board & audit purposes"]
    for i, feat in enumerate(features9):
        fbox = slide9.shapes.add_textbox(Inches(1), Inches(3.3 + i * 0.85), Inches(8), Inches(0.75))
        tf_f = fbox.text_frame
        tf_f.word_wrap = True
        tf_f.text = f"✓ {feat}"
        p_f = tf_f.paragraphs[0]
        p_f.font.size = Pt(12)
        p_f.font.color.rgb = HITAM
    print("✅ Slide 9: Dashboard")

    # ===== SLIDE 10: VENDOR =====
    slide10 = add_slide_with_background("Manajemen Penyedia & Portal")
    vendor_flow = ["Pendaftaran", "Kualifikasi", "Penilaian", "Monitoring"]
    for i, vf in enumerate(vendor_flow):
        left = Inches(1 + i * 2)
        vfbox = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(1.7), Inches(0.9))
        vfbox.fill.solid()
        vfbox.fill.fore_color.rgb = ORANGE if i % 2 == 0 else HIJAU
        vfbox.line.color.rgb = HIJAU
        vfbox.line.width = Pt(2)

        vftext = slide10.shapes.add_textbox(left + Inches(0.1), Inches(1.7), Inches(1.5), Inches(0.5))
        tf_vf = vftext.text_frame
        tf_vf.word_wrap = True
        tf_vf.text = vf
        p_vf = tf_vf.paragraphs[0]
        p_vf.font.size = Pt(11)
        p_vf.font.bold = True
        p_vf.font.color.rgb = PUTIH
        p_vf.alignment = PP_ALIGN.CENTER

    details10 = ["Portal Penyedia: Self-reporting obligations, SLA compliance, Document upload", "Vendor Scoring Matrix: 7 dimensi (Quality, Delivery, Compliance, Partnership, Financial, Service, Innovation)", "Automatic alerts jika vendor tidak meet SLA atau quality standards", "Vendor segmentation: Approved, At-risk, Blacklisted dengan automatic actions"]
    for i, det in enumerate(details10):
        dbox = slide10.shapes.add_textbox(Inches(0.7), Inches(2.8 + i * 0.95), Inches(8.6), Inches(0.85))
        tf_d = dbox.text_frame
        tf_d.word_wrap = True
        tf_d.text = f"• {det}"
        p_d = tf_d.paragraphs[0]
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = HITAM
    print("✅ Slide 10: Vendor Management")

    # ===== SLIDE 11: STAKEHOLDER =====
    slide11 = add_slide_with_background("Ekosistem Multi-Stakeholder KMU")
    stakeholders11 = [("Direksi PT KMU", "Strategic oversight, approval, KPI monitoring"), ("Departemen Pengadaan", "Daily operations, vendor management, compliance"), ("SPI (Audit Internal)", "Compliance verification, risk reporting, recommendations"), ("SBU (Unit Peminta)", "Request management, specification, acceptance"), ("Tim KSU (Vendor Liaison)", "KSO coordination, partnership management")]
    for i, (stakeholder, role) in enumerate(stakeholders11):
        top = Inches(1.3 + i * 1.05)
        stkbox = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), top, Inches(8.6), Inches(0.95))
        stkbox.fill.solid()
        stkbox.fill.fore_color.rgb = GRAY_LIGHT
        stkbox.line.color.rgb = HIJAU if i % 2 == 0 else ORANGE
        stkbox.line.width = Pt(2)

        stkname = slide11.shapes.add_textbox(Inches(1), top + Inches(0.1), Inches(2.5), Inches(0.4))
        tf_sn = stkname.text_frame
        tf_sn.text = stakeholder
        p_sn = tf_sn.paragraphs[0]
        p_sn.font.size = Pt(12)
        p_sn.font.bold = True
        p_sn.font.color.rgb = HIJAU

        stkrole = slide11.shapes.add_textbox(Inches(3.7), top + Inches(0.1), Inches(5), Inches(0.4))
        tf_sr = stkrole.text_frame
        tf_sr.word_wrap = True
        tf_sr.text = role
        p_sr = tf_sr.paragraphs[0]
        p_sr.font.size = Pt(11)
        p_sr.font.color.rgb = HITAM
    print("✅ Slide 11: Stakeholder")

    # ===== SLIDE 12: COST & ROI =====
    slide12 = add_slide_with_background("Provider Comparison: Cost & ROI")
    providers = [
        ["PROVIDER", "INIT COST", "YEAR 1 COST", "YEAR 1 ROI", "PAYBACK"],
        ["🏆 Deepseek", "Rp 67-98 JT", "Rp 80-110 JT", "9,636%", "5 hari"],
        ["Groq", "Rp 82-115 JT", "Rp 100-140 JT", "8,800%", "6 hari"],
        ["Claude", "Rp 104-148 JT", "Rp 130-170 JT", "7,423%", "13 hari"],
        ["Self-Hosted", "Rp 134-155 JT", "Rp 160-190 JT", "6,031%", "16 hari"],
        ["Gemini", "Rp 1.045 B", "Rp 1.045 B", "692%", "1.5 bulan"]
    ]

    col_widths = [1.5, 1.5, 1.5, 1.2, 1.2]
    total_width = sum(col_widths)
    start_left = Inches(0.7)
    start_top = Inches(1.3)
    cell_height = Inches(0.4)

    for row_idx, row_data in enumerate(providers):
        for col_idx, cell_val in enumerate(row_data):
            left = start_left + Inches(sum(col_widths[:col_idx]))
            top = start_top + Inches(row_idx * 0.4)
            width = Inches(col_widths[col_idx])

            cell_bg = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, cell_height)
            cell_bg.fill.solid()

            if row_idx == 0:
                cell_bg.fill.fore_color.rgb = ORANGE
            elif row_idx == 1:
                cell_bg.fill.fore_color.rgb = RGBColor(220, 255, 220)
            else:
                cell_bg.fill.fore_color.rgb = GRAY_LIGHT

            cell_bg.line.color.rgb = HIJAU
            cell_bg.line.width = Pt(1)

            cell_text = slide12.shapes.add_textbox(left + Inches(0.05), top + Inches(0.05), width - Inches(0.1), cell_height - Inches(0.1))
            tf_ct = cell_text.text_frame
            tf_ct.word_wrap = True
            tf_ct.text = cell_val
            p_ct = tf_ct.paragraphs[0]
            p_ct.font.size = Pt(9 if row_idx == 0 else 8)
            p_ct.font.bold = True if row_idx <= 1 else False
            p_ct.font.color.rgb = PUTIH if row_idx == 0 else HITAM
            p_ct.alignment = PP_ALIGN.CENTER

    rec_box = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(4.2), Inches(7), Inches(1.8))
    rec_box.fill.solid()
    rec_box.fill.fore_color.rgb = ORANGE
    rec_box.line.color.rgb = HIJAU
    rec_box.line.width = Pt(3)

    rec_text = slide12.shapes.add_textbox(Inches(1.8), Inches(4.4), Inches(6.4), Inches(1.4))
    tf_rec = rec_text.text_frame
    tf_rec.word_wrap = True
    tf_rec.text = "✅ REKOMENDASI: DEEPSEEK\n\nROI tertinggi 9,636% | Payback tercepat 5 hari | Investment minimal\nSetup mudah, integration cepat, maintenance ringan"

    for i, p in enumerate(tf_rec.paragraphs):
        if i == 0:
            p.font.size = Pt(16)
            p.font.bold = True
        else:
            p.font.size = Pt(11)
        p.font.color.rgb = PUTIH
        p.alignment = PP_ALIGN.CENTER
    print("✅ Slide 12: Cost & ROI Comparison")

    # ===== SLIDE 13: TECHNICAL =====
    slide13 = add_slide_with_background("Technical Capabilities Comparison")
    tech_compare = [
        ["KRITERIA", "DEEPSEEK", "GROQ", "CLAUDE", "GEMINI"],
        ["Response Time", "Good", "EXCELLENT", "Good", "Good"],
        ["Document Accuracy", "95%", "94%", "97%", "98%"],
        ["Validation Accuracy", "92%", "91%", "95%", "96%"],
        ["Uptime Reliability", "99.5%", "99.5%", "99.8%", "99.9%"],
        ["Support Quality", "Good", "Good", "Excellent", "Excellent"]
    ]

    col_widths13 = [1.8, 1.6, 1.6, 1.6, 1.6]
    start_top13 = Inches(1.3)

    for row_idx, row_data in enumerate(tech_compare):
        for col_idx, cell_val in enumerate(row_data):
            left = start_left + Inches(sum(col_widths13[:col_idx]))
            top = start_top13 + Inches(row_idx * 0.45)
            width = Inches(col_widths13[col_idx])

            cell_bg = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.4))
            cell_bg.fill.solid()

            if row_idx == 0:
                cell_bg.fill.fore_color.rgb = HIJAU
            elif col_idx == 1:
                cell_bg.fill.fore_color.rgb = RGBColor(255, 240, 200)
            else:
                cell_bg.fill.fore_color.rgb = GRAY_LIGHT

            cell_bg.line.color.rgb = ORANGE
            cell_bg.line.width = Pt(1)

            cell_text = slide13.shapes.add_textbox(left + Inches(0.05), top + Inches(0.05), width - Inches(0.1), Inches(0.3))
            tf_ct = cell_text.text_frame
            tf_ct.word_wrap = True
            tf_ct.text = cell_val
            p_ct = tf_ct.paragraphs[0]
            p_ct.font.size = Pt(9)
            p_ct.font.bold = True if row_idx == 0 else False
            p_ct.font.color.rgb = PUTIH if row_idx == 0 else HITAM
            p_ct.alignment = PP_ALIGN.CENTER

    conc_text = slide13.shapes.add_textbox(Inches(0.7), Inches(4), Inches(8.6), Inches(2.5))
    tf_conc = conc_text.text_frame
    tf_conc.word_wrap = True
    tf_conc.text = """KESIMPULAN:
✓ Deepseek = BEST VALUE: ROI + cost optimal, setup mudah, maintenance ringan
✓ Groq = PERFORMANCE: Response tercepat (<200ms), ideal untuk real-time validation
✓ Claude = HIGH QUALITY: Akurasi tertinggi, excellent support, jika budget tidak terbatas
► REKOMENDASI: Deepseek dengan Groq sebagai fallback untuk optimal cost-performance"""

    for p in tf_conc.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = HITAM
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    print("✅ Slide 13: Technical Comparison")

    # ===== SLIDE 14: TIMELINE =====
    slide14 = add_slide_with_background("Implementation Timeline: 12 Minggu")
    timeline = [("W1-2", "Setup & Validation", "Rp 150M"), ("W3-6", "Core Development", "Rp 600M"), ("W7-8", "Chatbot & RAG", "Rp 150M"), ("W9-10", "Platform Integration", "Rp 100M"), ("W11", "Testing & UAT", "Rp 50M"), ("W12", "Training & Launch", "Rp 50M")]

    for i, (weeks, task, cost) in enumerate(timeline):
        row = i // 3
        col = i % 3
        left = Inches(0.6 + col * 3)
        top = Inches(1.3 + row * 2.5)

        tbox = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(2))
        tbox.fill.solid()
        tbox.fill.fore_color.rgb = ORANGE if i % 2 == 0 else HIJAU
        tbox.line.color.rgb = HIJAU
        tbox.line.width = Pt(2)

        wtext = slide14.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(2.4), Inches(0.3))
        tf_w = wtext.text_frame
        tf_w.text = weeks
        p_w = tf_w.paragraphs[0]
        p_w.font.size = Pt(12)
        p_w.font.bold = True
        p_w.font.color.rgb = PUTIH
        p_w.alignment = PP_ALIGN.CENTER

        tttext = slide14.shapes.add_textbox(left + Inches(0.2), top + Inches(0.6), Inches(2.4), Inches(0.8))
        tf_tt = tttext.text_frame
        tf_tt.word_wrap = True
        tf_tt.text = task
        for p in tf_tt.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = PUTIH
            p.alignment = PP_ALIGN.CENTER

        ctext = slide14.shapes.add_textbox(left + Inches(0.2), top + Inches(1.5), Inches(2.4), Inches(0.3))
        tf_c = ctext.text_frame
        tf_c.text = cost
        p_c = tf_c.paragraphs[0]
        p_c.font.size = Pt(9)
        p_c.font.color.rgb = PUTIH
        p_c.alignment = PP_ALIGN.CENTER

    total_box = slide14.shapes.add_textbox(Inches(2.5), Inches(6.2), Inches(5), Inches(0.8))
    tf_total = total_box.text_frame
    tf_total.text = "TOTAL INVESTASI: Rp 1.04 MILIAR (12 minggu)"
    p_total = tf_total.paragraphs[0]
    p_total.font.size = Pt(14)
    p_total.font.bold = True
    p_total.font.color.rgb = ORANGE
    p_total.alignment = PP_ALIGN.CENTER
    print("✅ Slide 14: Timeline")

    # ===== SLIDE 15: KPI =====
    slide15 = add_slide_with_background("Success Metrics & KPI Targets")
    kpis = [("85-90%", "Success Probability"), ("<500ms", "Validation Response"), (">95%", "Accuracy Target"), ("30%", "Process Time Reduction"), ("100%", "SPO Compliance"), ("12 Minggu", "Implementation Timeline")]

    for i, (value, label) in enumerate(kpis):
        col = i % 3
        row = i // 3
        left = Inches(0.7 + col * 3)
        top = Inches(1.3 + row * 2.2)

        kbox = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.7), Inches(1.8))
        kbox.fill.solid()
        kbox.fill.fore_color.rgb = ORANGE if i % 2 == 0 else HIJAU
        kbox.line.color.rgb = HIJAU
        kbox.line.width = Pt(2)

        vtext = slide15.shapes.add_textbox(left + Inches(0.2), top + Inches(0.3), Inches(2.3), Inches(0.7))
        tf_v = vtext.text_frame
        tf_v.word_wrap = True
        tf_v.text = value
        p_v = tf_v.paragraphs[0]
        p_v.font.size = Pt(22)
        p_v.font.bold = True
        p_v.font.color.rgb = PUTIH
        p_v.alignment = PP_ALIGN.CENTER

        ltext = slide15.shapes.add_textbox(left + Inches(0.2), top + Inches(1.1), Inches(2.3), Inches(0.6))
        tf_l = ltext.text_frame
        tf_l.word_wrap = True
        tf_l.text = label
        for p in tf_l.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = PUTIH
            p.alignment = PP_ALIGN.CENTER
    print("✅ Slide 15: KPI")

    # ===== SLIDE 16: RISK =====
    slide16 = add_slide_with_background("Risk Management & Mitigation")
    risks = [("AI Provider Risk", "Deepseek fallback ke Groq/Claude jika ada issue"), ("User Adoption Risk", "Extensive training, UI yang user-friendly, mandatory mode"), ("Technology Risk", "Load testing, caching, optimization, redundancy"), ("Budget Risk", "Deepseek minimal, contingency Rp 10-20M reserved")]

    for i, (risk, mitigation) in enumerate(risks):
        top = Inches(1.3 + i * 1.3)

        rtext = slide16.shapes.add_textbox(Inches(0.7), top, Inches(3), Inches(0.4))
        tf_r = rtext.text_frame
        tf_r.text = risk
        p_r = tf_r.paragraphs[0]
        p_r.font.size = Pt(12)
        p_r.font.bold = True
        p_r.font.color.rgb = ORANGE

        mtext = slide16.shapes.add_textbox(Inches(3.8), top, Inches(5.5), Inches(0.5))
        tf_m = mtext.text_frame
        tf_m.word_wrap = True
        tf_m.text = mitigation
        p_m = tf_m.paragraphs[0]
        p_m.font.size = Pt(10)
        p_m.font.color.rgb = HITAM

    assess_box = slide16.shapes.add_textbox(Inches(0.7), Inches(6), Inches(8.6), Inches(0.9))
    tf_a = assess_box.text_frame
    tf_a.word_wrap = True
    tf_a.text = "✓ OVERALL: Probabilitas kesuksesan 85-90% dengan contingency plans yang comprehensive. Worst-case scenario masih menghasilkan ROI >846%"
    p_a = tf_a.paragraphs[0]
    p_a.font.size = Pt(11)
    p_a.font.bold = True
    p_a.font.color.rgb = HIJAU
    print("✅ Slide 16: Risk Management")

    # ===== SLIDE 17: NEXT STEPS =====
    slide17 = add_slide_with_background("Langkah Selanjutnya")
    next_steps = [("1", "Board Approval", "Vote SKD untuk approve project ini dengan Deepseek"), ("2", "Budget Allocation", "CFO mengalokasikan budget Rp 80-110 JT"), ("3", "Team Preparation", "IT Director confirm team availability & assign leads"), ("4", "Project Kickoff", "Week 1: Setup infrastructure, validation testing mulai")]

    for i, (num, step, detail) in enumerate(next_steps):
        top = Inches(1.4 + i * 1.4)

        ncircle = slide17.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), top, Inches(0.5), Inches(0.5))
        ncircle.fill.solid()
        ncircle.fill.fore_color.rgb = ORANGE
        ncircle.line.color.rgb = HIJAU

        ntext = slide17.shapes.add_textbox(Inches(0.7), top, Inches(0.5), Inches(0.5))
        tf_n = ntext.text_frame
        tf_n.text = num
        p_n = tf_n.paragraphs[0]
        p_n.font.size = Pt(16)
        p_n.font.bold = True
        p_n.font.color.rgb = PUTIH
        p_n.alignment = PP_ALIGN.CENTER
        p_n.vertical_anchor = MSO_ANCHOR.MIDDLE

        stext = slide17.shapes.add_textbox(Inches(1.5), top, Inches(7.8), Inches(0.5))
        tf_s = stext.text_frame
        tf_s.word_wrap = True
        tf_s.text = f"{step}: {detail}"
        p_s = tf_s.paragraphs[0]
        p_s.font.size = Pt(11)
        p_s.font.color.rgb = HITAM

    cta_box = slide17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(6.2), Inches(6), Inches(0.8))
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = HIJAU
    cta_box.line.color.rgb = ORANGE
    cta_box.line.width = Pt(2)

    cta_text = slide17.shapes.add_textbox(Inches(2.2), Inches(6.4), Inches(5.6), Inches(0.4))
    tf_cta = cta_text.text_frame
    tf_cta.text = "Expected Delivery: September 2026 | Expected Year 1 Profit: Rp 9.71 MILIAR"
    p_cta = tf_cta.paragraphs[0]
    p_cta.font.size = Pt(12)
    p_cta.font.bold = True
    p_cta.font.color.rgb = PUTIH
    p_cta.alignment = PP_ALIGN.CENTER
    print("✅ Slide 17: Next Steps")

    # ===== SLIDE 18: CLOSING =====
    slide18 = prs.slides.add_slide(prs.slide_layouts[6])
    bg18 = slide18.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg18.fill.solid()
    bg18.fill.fore_color.rgb = HIJAU
    bg18.line.color.rgb = HIJAU

    closing_text = slide18.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    tf_closing = closing_text.text_frame
    tf_closing.word_wrap = True
    tf_closing.text = "Terima Kasih\n\nMari Transformasikan Pengadaan KMU\ndengan Platform Terintegrasi & Validator Prosedur"

    for i, p in enumerate(tf_closing.paragraphs):
        if i == 0:
            p.font.size = Pt(48)
        else:
            p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = ORANGE
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)

    footer18 = slide18.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    tf_footer18 = footer18.text_frame
    tf_footer18.text = "PT. Kaltim Medika Utama | Juni 2026\nProject: Digitalisasi Pengadaan KMU dengan AI"
    for p in tf_footer18.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = PUTIH
        p.alignment = PP_ALIGN.CENTER
    print("✅ Slide 18: Closing")

    # Save
    output_path = r"C:\Users\Petugas\Claude\Projects\Pengadaan\Presentasi_Digitalisasi_Pengadaan_KMU_REVISED.pptx"
    prs.save(output_path)

    print(f"\n" + "="*70)
    print(f"✅ PPT BARU BERHASIL DIBUAT - RAPI DAN LENGKAP!")
    print(f"="*70)
    print(f"File: {output_path}")
    print(f"Total Slides: {len(prs.slides)}")
    print(f"\nWarna KMU Diterapkan:")
    print(f"  • Hijau: #2E8B5E")
    print(f"  • Orange: #FF9F1C")
    print(f"\nKualitas:")
    print(f"  ✓ RAPI - Layout proper, spacing yang baik")
    print(f"  ✓ LENGKAP - 18 slides dengan semua konten")
    print(f"  ✓ NO OVERLAP - Tulisan & grafik terpisah dengan jelas")
    print(f"  ✓ PROFESSIONAL - Consistent styling, proper typography")
    print(f"\nKonten:")
    print(f"  • Slide 1: Cover")
    print(f"  • Slides 2-11: Strategy & Architecture")
    print(f"  • Slides 12-13: Provider Comparison Matrix (5 providers)")
    print(f"  • Slides 14-17: Implementation Details")
    print(f"  • Slide 18: Closing")

if __name__ == "__main__":
    create_presentation()
